# 文件说明：后端知识库服务，负责资料解析、保存和检索；当前文件：document parser
"""PDF / PPTX 文本提取与语义分段。"""

from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path

from pypdf import PdfReader
from pptx import Presentation


@dataclass
class ParsedSection:
    title: str | None
    content: str
    metadata: dict


def _normalize_text(text: str) -> str:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+\n", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _split_long_text(
    text: str,
    *,
    title: str | None,
    metadata: dict,
    target_chars: int = 750,
    max_chars: int = 1000,
) -> list[ParsedSection]:
    text = _normalize_text(text)
    if not text:
        return []

    paragraphs = [p.strip() for p in re.split(r"\n{2,}", text) if p.strip()]
    if not paragraphs:
        return []

    chunks: list[str] = []
    current = ""
    for paragraph in paragraphs:
        if len(paragraph) > max_chars:
            if current:
                chunks.append(current.strip())
                current = ""
            for i in range(0, len(paragraph), target_chars):
                chunks.append(paragraph[i : i + target_chars].strip())
            continue
        candidate = f"{current}\n\n{paragraph}".strip() if current else paragraph
        if len(candidate) <= target_chars:
            current = candidate
        else:
            if current:
                chunks.append(current.strip())
            current = paragraph
    if current:
        chunks.append(current.strip())

    return [
        ParsedSection(title=title, content=chunk, metadata=dict(metadata))
        for chunk in chunks
        if chunk
    ]


def parse_pdf(file_path: Path) -> list[ParsedSection]:
    reader = PdfReader(str(file_path))
    sections: list[ParsedSection] = []
    for page_index, page in enumerate(reader.pages, start=1):
        raw = page.extract_text() or ""
        text = _normalize_text(raw)
        if not text:
            continue
        title = f"第 {page_index} 页"
        sections.extend(
            _split_long_text(
                text,
                title=title,
                metadata={"page": page_index, "source_type": "pdf"},
            )
        )
    return sections


def _slide_text_parts(slide) -> tuple[str | None, str]:
    title = None
    body_parts: list[str] = []
    for shape in slide.shapes:
        if not hasattr(shape, "text"):
            continue
        text = _normalize_text(shape.text or "")
        if not text:
            continue
        if title is None and getattr(shape, "is_placeholder", False):
            try:
                if shape.placeholder_format.type == 1:  # TITLE
                    title = text
                    continue
            except Exception:
                pass
        body_parts.append(text)
    if title is None and body_parts:
        first = body_parts[0]
        if len(first) <= 80:
            title = first
            body_parts = body_parts[1:]
    body = _normalize_text("\n".join(body_parts))
    return title, body


def parse_pptx(file_path: Path) -> list[ParsedSection]:
    presentation = Presentation(str(file_path))
    sections: list[ParsedSection] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        title, body = _slide_text_parts(slide)
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = _normalize_text(slide.notes_slide.notes_text_frame.text or "")

        combined_parts = [part for part in (body, notes) if part]
        combined = _normalize_text("\n\n".join(combined_parts))
        if not combined:
            continue

        section_title = title or f"幻灯片 {slide_index}"
        metadata = {
            "slide": slide_index,
            "source_type": "pptx",
            "has_notes": bool(notes),
        }
        if title:
            metadata["slide_title"] = title
        sections.extend(
            _split_long_text(
                combined,
                title=section_title,
                metadata=metadata,
            )
        )
    return sections


def parse_document(file_path: Path) -> list[ParsedSection]:
    suffix = file_path.suffix.lower()
    if suffix == ".pdf":
        return parse_pdf(file_path)
    if suffix in {".pptx", ".ppt"}:
        if suffix == ".ppt":
            raise ValueError("暂不支持旧版 .ppt 格式，请转换为 .pptx")
        return parse_pptx(file_path)
    raise ValueError(f"不支持的文件类型: {suffix}")
