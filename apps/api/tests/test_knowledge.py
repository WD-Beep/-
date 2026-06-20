"""知识库与话术推荐基础测试。"""

from __future__ import annotations

import asyncio
from pathlib import Path

import pytest
from sqlalchemy import delete

from app.db.session import async_session_factory
from app.deps.tenant import TenantContext
from app.models.knowledge import KnowledgeBase, KnowledgeChunk, KnowledgeDocument
from app.services.knowledge.document_parser import parse_document


def test_knowledge_routes_reject_all_products_scope():
    async def _run() -> None:
        from httpx import ASGITransport, AsyncClient

        from app.main import app

        transport = ASGITransport(app=app)
        async with AsyncClient(transport=transport, base_url="http://test") as client:
            response = await client.get(
                "/api/knowledge-bases",
                headers={"X-User-Id": "1", "X-Product-Id": "0"},
            )
            assert response.status_code == 400

    asyncio.run(_run())


def test_document_parser_pdf_and_pptx_smoke(tmp_path: Path):
    pdf_path = tmp_path / "sample.pdf"
    pptx_path = tmp_path / "sample.pptx"

    from pypdf import PdfWriter
    from pptx import Presentation

    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    with pdf_path.open("wb") as handle:
        writer.write(handle)

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "ScandiHome"
    slide.placeholders[1].text = "北欧简约家居品牌视觉升级"
    presentation.save(str(pptx_path))

    pdf_sections = parse_document(pdf_path)
    pptx_sections = parse_document(pptx_path)
    assert isinstance(pdf_sections, list)
    assert len(pptx_sections) >= 1
    assert "ScandiHome" in pptx_sections[0].content or "北欧" in pptx_sections[0].content


def test_knowledge_search_keyword_match():
    async def _run() -> None:
        from app.services.knowledge.knowledge_service import KnowledgeService
        from app.services.knowledge.search_service import KnowledgeSearchService

        async with async_session_factory() as db:
            await db.execute(delete(KnowledgeChunk))
            await db.execute(delete(KnowledgeDocument))
            await db.execute(delete(KnowledgeBase))
            await db.commit()

            ctx = TenantContext(user_id=1, product_id=1, workspace_id=1, is_admin=True)
            base = await KnowledgeService.get_or_create_default_base(db, ctx=ctx, product_id=1)
            doc = KnowledgeDocument(
                knowledge_base_id=base.id,
                workspace_id=1,
                product_id=1,
                file_name="brand.pdf",
                file_type="pdf",
                status="ready",
            )
            db.add(doc)
            await db.flush()
            db.add(
                KnowledgeChunk(
                    document_id=doc.id,
                    knowledge_base_id=base.id,
                    workspace_id=1,
                    product_id=1,
                    chunk_index=0,
                    title="品牌定位",
                    content="ScandiHome 专注北欧简约家居，强调自然材质与克制配色。",
                    chunk_metadata={"page": 1},
                )
            )
            await db.commit()

            hits = await KnowledgeSearchService.search(
                db, product_id=1, query="北欧 家居 视觉", limit=5
            )
            assert len(hits) >= 1
            assert hits[0].document_name == "brand.pdf"

    asyncio.run(_run())
